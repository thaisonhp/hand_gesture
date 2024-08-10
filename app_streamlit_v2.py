import cv2
import streamlit as st
import mediapipe as mp
import numpy as np
import torch
import requests
from pptx import Presentation
from PIL import Image
import io
from tqdm import tqdm
from scipy.spatial.distance import cdist
from scipy.ndimage import zoom, median_filter as medfilt
import random

class Config:
    def __init__(self):
        self.frame_l = 32  # the length of frames
        self.joint_n = 22  # the number of joints
        self.joint_d = 3  # the dimension of joints
        self.clc_num = 14  # the number of class
        self.feat_d = 231
        self.filters = 64

def pad_arrays(data):
    max_length = max(max(arr.shape[0] for arr in frame) for frame in data)
    return np.array([[np.pad(arr, (0, max_length - arr.shape[0]), 'constant') for arr in frame] for frame in data])

def zoom_and_filter(p, target_l, joints_num, joints_dim):
    l = p.shape[0]
    p_new = np.empty(shape=[target_l, joints_num, joints_dim])
    for m in range(joints_num):
        for n in range(joints_dim):
            p[:, m, n] = medfilt(p[:, m, n], 3)
            p_new[:, m, n] = zoom(p[:, m, n], target_l / l)[:target_l]
    return p_new

def sampling_frame(p, C):
    full_l = p.shape[0]
    if random.uniform(0, 1) < 0.5:  # alignment sampling
        valid_l = np.round(np.random.uniform(0.9, 1) * full_l)
        s = random.randint(0, full_l - int(valid_l))
        e = s + valid_l
        p = p[int(s):int(e), :, :]
    else:  # without alignment sampling
        valid_l = np.round(np.random.uniform(0.9, 1) * full_l)
        index = np.sort(np.random.choice(range(full_l), int(valid_l), replace=False))
        p = p[index, :, :]
    p = zoom_and_filter(p, C.frame_l, C.joint_n, C.joint_d)
    return p

def norm_train(p):
    p[:, :, 0] -= np.mean(p[:, :, 0])
    p[:, :, 1] -= np.mean(p[:, :, 1])
    p[:, :, 2] -= np.mean(p[:, :, 2])
    return p

def norm_train2d(p):
    p[:, :, 0] -= np.mean(p[:, :, 0])
    p[:, :, 1] -= np.mean(p[:, :, 1])
    return p

def get_CG(p, C):
    M = []
    iu = np.triu_indices(C.joint_n, 1, C.joint_n)
    for f in range(C.frame_l):
        d_m = cdist(p[f], np.concatenate([p[f], np.zeros([1, C.joint_d])]), 'euclidean')
        d_m = d_m[iu]
        M.append(d_m)
    M = np.stack(M)
    return M

def data_generator_rt(T, C):
    X_0 = []
    X_1 = []
    T = np.expand_dims(T, axis=0)
    for i in tqdm(range(len(T))):
        p = np.copy(T[i])
        p = zoom_and_filter(p, target_l=C.frame_l, joints_num=C.joint_n, joints_dim=C.joint_d)
        M = get_CG(p, C)
        X_0.append(M)
        p = norm_train2d(p)
        X_1.append(p)
    X_0 = np.stack(X_0)
    X_1 = np.stack(X_1)
    return X_0, X_1

def upload_pptx():
    """Tải lên file PPTX và trả về đối tượng Presentation"""
    uploaded_file = st.file_uploader("Tải lên file PPTX", type="pptx")
    if uploaded_file is not None:
        return Presentation(uploaded_file)
    return None

def show_slide_content(presentation, slide_index):
    """Hiển thị nội dung và hình ảnh của slide tại slide_index"""
    slide = presentation.slides[slide_index]
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            st.write(shape.text)
        if hasattr(shape, "image"):
            img_stream = io.BytesIO(shape.image.blob)
            img = Image.open(img_stream)
            st.image(img)

def main():
    st.title("Webcam Live Feed with Hand Tracking")

    # Điều khiển YouTube
    video_url = st.text_input("Nhập URL video YouTube:")
    if video_url:
        st.video(video_url)
    else:
        st.write("Vui lòng nhập URL video YouTube.")

    # Điều khiển PPTX
    presentation = upload_pptx()
    if presentation is not None:
        slide_count = len(presentation.slides)
        current_slide_index = st.session_state.get('current_slide_index', 0)
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button('Previous Slide'):
                if current_slide_index > 0:
                    current_slide_index -= 1
                    st.session_state['current_slide_index'] = current_slide_index
        
        with col2:
            if st.button('Next Slide'):
                if current_slide_index < slide_count - 1:
                    current_slide_index += 1
                    st.session_state['current_slide_index'] = current_slide_index
        
        show_slide_content(presentation, current_slide_index)
        st.write(f"Slide {current_slide_index + 1} of {slide_count}")

    run = st.checkbox('Run')
    FRAME_WINDOW = st.image([])
    C = Config()
    mp_hands = mp.solutions.hands
    hands = mp_hands.Hands(static_image_mode=False, max_num_hands=2, min_detection_confidence=0.5, min_tracking_confidence=0.5)
    mp_drawing = mp.solutions.drawing_utils

    sequence = []
    predictions = []
    camera = cv2.VideoCapture(0)

    while run:
        ret, frame = camera.read()
        if not ret:
            break

        frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        results = hands.process(frame)

        if results.multi_hand_landmarks:
            for hand_landmarks in results.multi_hand_landmarks:
                mp_drawing.draw_landmarks(frame, hand_landmarks, mp_hands.HAND_CONNECTIONS)
                keypoint = np.array([[res.x, res.y, res.z] for res in hand_landmarks.landmark]).flatten()
                keypoint = np.array_split(keypoint, 22)
                sequence.append(keypoint)
                sequence = sequence[-32:]

            if len(sequence) == 32:
                sequence_np = np.array(sequence, dtype=object)
                sequence_np = pad_arrays(sequence_np)
                X_test_rt_1, X_test_rt_2 = data_generator_rt(sequence_np[-32:], C)

                X_test_rt_1 = torch.from_numpy(X_test_rt_1).type(torch.FloatTensor)
                X_test_rt_2 = torch.from_numpy(X_test_rt_2).type(torch.FloatTensor)
                X_test_rt_1_list = X_test_rt_1.tolist()
                X_test_rt_2_list = X_test_rt_2.tolist()

                try:
                    response = requests.post("http://127.0.0.1:8000/predict", json={
                        "X_test_rt_1": X_test_rt_1_list,
                        "X_test_rt_2": X_test_rt_2_list
                    })
                    
                    if response.status_code == 200:
                        prediction = response.json().get("prediction", "Unknown")
                        predictions.append(prediction)
                        if predictions[-1] != predictions[-2]:
                            # call your YouTube control function here
                            st.write(f"Prediction: {predictions[-1]}")
                    else:
                        st.write("Error:", response.text)
                except Exception as e:
                    st.write(f"Request failed: {str(e)}")
        FRAME_WINDOW.image(frame)
    else:
        st.write('Stopped')

    camera.release()

if __name__ == "__main__":
    main()
